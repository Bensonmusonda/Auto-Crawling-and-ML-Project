import sys
import os

# Ensure python-pptx is installed
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("Installing python-pptx library...")
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    
    # Set slide dimensions to widescreen (16:9)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    blank_slide_layout = prs.slide_layouts[6]
    
    # Themes / Colors (Monochrome Zinc/Dark Theme with Blue accents)
    COLOR_BG = RGBColor(9, 9, 11)          # #09090b
    COLOR_CARD = RGBColor(24, 24, 27)      # #18181b
    COLOR_BORDER = RGBColor(39, 39, 42)    # #27272a
    COLOR_TEXT = RGBColor(244, 244, 245)    # #f4f4f5
    COLOR_MUTED = RGBColor(161, 161, 170)  # #a1a1aa
    COLOR_ACCENT = RGBColor(59, 130, 246)  # #3b82f6
    COLOR_SUCCESS = RGBColor(34, 197, 94)  # #22c55e

    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_BG

    def add_header(slide, tag_text, title_text):
        # Tag text box
        tag_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf_tag = tag_box.text_frame
        tf_tag.word_wrap = True
        p_tag = tf_tag.paragraphs[0]
        p_tag.text = tag_text.upper()
        p_tag.font.name = "Arial"
        p_tag.font.size = Pt(11)
        p_tag.font.bold = True
        p_tag.font.color.rgb = COLOR_ACCENT
        
        # Title text box
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.name = "Arial"
        p_title.font.size = Pt(32)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_TEXT

    # ── Slide 1: Title Slide ──────────────────────────────────────────
    slide1 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide1)
    
    # Main container shape
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.33), Inches(4.0))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    
    # Tag
    p_tag = tf1.paragraphs[0]
    p_tag.text = "CAPSTONE DEFENCE PRESENTATION"
    p_tag.font.name = "Arial"
    p_tag.font.size = Pt(12)
    p_tag.font.bold = True
    p_tag.font.color.rgb = COLOR_ACCENT
    p_tag.alignment = PP_ALIGN.CENTER
    
    # Project Title
    p_title = tf1.add_paragraph()
    p_title.text = "\nAN AUTOMATED DATA ACQUISITION AND\nMACHINE LEARNING MODEL GENERATION SYSTEM"
    p_title.font.name = "Arial"
    p_title.font.size = Pt(34)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_TEXT
    p_title.alignment = PP_ALIGN.CENTER
    
    # Presenter Details
    p_pres = tf1.add_paragraph()
    p_pres.text = "\nPresented by: Benson Musonda (Student ID: 202204757)"
    p_pres.font.name = "Arial"
    p_pres.font.size = Pt(15)
    p_pres.font.bold = True
    p_pres.font.color.rgb = COLOR_MUTED
    p_pres.alignment = PP_ALIGN.CENTER
    
    p_sup = tf1.add_paragraph()
    p_sup.text = "Supervisor: Dr. Sinyinda Muwanei"
    p_sup.font.name = "Arial"
    p_sup.font.size = Pt(13)
    p_sup.font.color.rgb = COLOR_MUTED
    p_sup.alignment = PP_ALIGN.CENTER

    # ── Slide 2: Introduction ─────────────────────────────────────────
    slide2 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide2)
    add_header(slide2, "Background Context", "Introduction")
    
    # Highlight Box (Central theme quote)
    shape_h = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.7), Inches(1.6))
    shape_h.fill.solid()
    shape_h.fill.fore_color.rgb = COLOR_CARD
    shape_h.line.color.rgb = COLOR_ACCENT
    shape_h.line.width = Pt(3)
    
    tf_h = shape_h.text_frame
    tf_h.word_wrap = True
    tf_h.margin_left = Inches(0.3)
    tf_h.margin_right = Inches(0.3)
    p_h = tf_h.paragraphs[0]
    p_h.text = '"Central to any effective machine learning application is the quality and availability of its training data. Furthermore, for these models to be truly valuable and trustworthy, their entire development lifecycle must be fully reproducible."'
    p_h.font.name = "Arial"
    p_h.font.size = Pt(16)
    p_h.font.color.rgb = COLOR_TEXT
    p_h.alignment = PP_ALIGN.CENTER
    
    # Description Body Box
    desc_box = slide2.shapes.add_textbox(Inches(0.8), Inches(3.7), Inches(11.7), Inches(3.2))
    tf_desc = desc_box.text_frame
    tf_desc.word_wrap = True
    
    p_d1 = tf_desc.paragraphs[0]
    p_d1.text = "The Automated Data Acquisition and Machine Learning Model Generation System addresses a significant challenge: bridging the gap between disparate data collection and AI model development tools."
    p_d1.font.name = "Arial"
    p_d1.font.size = Pt(16)
    p_d1.font.color.rgb = COLOR_MUTED
    
    p_d2 = tf_desc.add_paragraph()
    p_d2.text = "\nMany domain experts struggle to leverage powerful web data acquisition or advanced machine learning due to technical complexities. This system provides a unified, user-friendly platform that seamlessly integrates data ingestion from diverse sources (including robust web scraping) with adaptive data preprocessing and advanced text feature engineering, leading to automated, yet customizable, machine learning model generation."
    p_d2.font.name = "Arial"
    p_d2.font.size = Pt(16)
    p_d2.font.color.rgb = COLOR_MUTED

    # ── Slide 3: Problem Statement ────────────────────────────────────
    slide3 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide3)
    add_header(slide3, "The Challenge", "Problem Statement")
    
    # Description
    desc_box3 = slide3.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(1.0))
    tf_desc3 = desc_box3.text_frame
    tf_desc3.word_wrap = True
    p_d3 = tf_desc3.paragraphs[0]
    p_d3.text = "The journey from raw data to a deployable machine learning model is often complex, fragmented, and technically demanding, creating barriers for domain experts:"
    p_d3.font.name = "Arial"
    p_d3.font.size = Pt(16)
    p_d3.font.color.rgb = COLOR_MUTED
    
    # 3-column cards
    cols = [
        {"title": "Fragmentation", "text": "Existing tools for data acquisition, preprocessing, feature engineering, and model training operate in silos, requiring manual integration and extensive coding."},
        {"title": "Technical Barrier", "text": "Leveraging powerful techniques like robust web scraping or Deep Learning for NLP remains inaccessible to users without deep programming expertise."},
        {"title": "Lack of Reproducibility", "text": "Ensuring that an entire ML pipeline, from data source to final model, can be reliably re-executed and verified is a persistent and complex challenge."}
    ]
    
    for i, c in enumerate(cols):
        left = Inches(0.8 + i * 4.0)
        top = Inches(2.7)
        width = Inches(3.7)
        height = Inches(4.0)
        
        card = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD
        card.line.color.rgb = COLOR_BORDER
        card.line.width = Pt(1)
        
        tf_card = card.text_frame
        tf_card.word_wrap = True
        tf_card.margin_left = Inches(0.2)
        tf_card.margin_right = Inches(0.2)
        tf_card.margin_top = Inches(0.2)
        
        p_ct = tf_card.paragraphs[0]
        p_ct.text = c["title"]
        p_ct.font.name = "Arial"
        p_ct.font.size = Pt(18)
        p_ct.font.bold = True
        p_ct.font.color.rgb = COLOR_ACCENT
        
        p_cx = tf_card.add_paragraph()
        p_cx.text = "\n" + c["text"]
        p_cx.font.name = "Arial"
        p_cx.font.size = Pt(14)
        p_cx.font.color.rgb = COLOR_MUTED

    # ── Slide 4: Aim ──────────────────────────────────────────────────
    slide4 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide4)
    add_header(slide4, "Core Target", "Aim: Unifying Data Acquisition & ML")
    
    # Highlight Box (Core Aim statement)
    shape_aim = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.6), Inches(11.7), Inches(1.3))
    shape_aim.fill.solid()
    shape_aim.fill.fore_color.rgb = COLOR_CARD
    shape_aim.line.color.rgb = COLOR_BORDER
    shape_aim.line.width = Pt(1)
    
    tf_aim = shape_aim.text_frame
    tf_aim.word_wrap = True
    tf_aim.margin_left = Inches(0.3)
    tf_aim.margin_right = Inches(0.3)
    p_aim = tf_aim.paragraphs[0]
    p_aim.text = 'This project aims to design, develop, and evaluate a unified, user-friendly Automated Data Acquisition and Machine Learning Model Generation System. By abstracting away technical complexities, the system empowers domain experts to efficiently leverage diverse data sources for robust, reproducible AI model development.'
    p_aim.font.name = "Arial"
    p_aim.font.size = Pt(14)
    p_aim.font.color.rgb = COLOR_TEXT
    
    # 4 Key Goals as 2x2 grid cards
    goals = [
        ("Automate ML Lifecycle", "Automate core aspects of the machine learning lifecycle for selected supervised learning models and sentiment analysis, from data acquisition to model deployment."),
        ("Adaptive Preprocessing", "Provide adaptive data preprocessing and sophisticated text feature engineering to effectively transform diverse, real-world data into high-quality features suitable for predictive modeling."),
        ("Robust Ingestion Capabilities", "Integrate robust dual-source data ingestion capabilities, supporting both structured file uploads and advanced hybrid web scraping."),
        ("Workflow Persistence", "Ensure comprehensive Workflow Persistence, guaranteeing complete reproducibility of the entire data-to-model pipeline for reliable and verifiable results.")
    ]
    
    for i, g in enumerate(goals):
        row = i // 2
        col = i % 2
        
        left = Inches(0.8 + col * 5.95)
        top = Inches(3.2 + row * 1.95)
        width = Inches(5.75)
        height = Inches(1.7)
        
        card = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD
        card.line.color.rgb = COLOR_BORDER
        card.line.width = Pt(1)
        
        tf_g = card.text_frame
        tf_g.word_wrap = True
        tf_g.margin_left = Inches(0.2)
        tf_g.margin_right = Inches(0.2)
        tf_g.margin_top = Inches(0.15)
        
        p_gt = tf_g.paragraphs[0]
        p_gt.text = f"{i+1:02d}. {g[0]}"
        p_gt.font.name = "Arial"
        p_gt.font.size = Pt(16)
        p_gt.font.bold = True
        p_gt.font.color.rgb = COLOR_ACCENT
        
        p_gd = tf_g.add_paragraph()
        p_gd.text = g[1]
        p_gd.font.name = "Arial"
        p_gd.font.size = Pt(12)
        p_gd.font.color.rgb = COLOR_MUTED

    # ── Slide 5: Objectives ───────────────────────────────────────────
    slide5 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide5)
    add_header(slide5, "Execution Roadmap", "Objectives")
    
    objectives = [
        "Develop a Hybrid Data Ingestion Engine (Files + Scrapy/Playwright).",
        "Implement an Interactive Schema Mapping Interface with real-time validation.",
        "Build an Adaptive Preprocessing & Feature Engineering Pipeline (Auto-cleaning + TF-IDF).",
        "Integrate an Automated & Customizable ML Core (AutoML + Hyperparameter tuning).",
        "Develop a Model Evaluation Dashboard (e.g. RMSE, Confusion Matrix, Accuracy).",
        "Implement Workflow Persistence for one-click retraining and reproducibility."
    ]
    
    for i, obj in enumerate(objectives):
        left = Inches(0.8)
        top = Inches(1.6 + i * 0.9)
        width = Inches(11.7)
        height = Inches(0.7)
        
        row_shape = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        row_shape.fill.solid()
        row_shape.fill.fore_color.rgb = COLOR_CARD
        row_shape.line.color.rgb = COLOR_BORDER
        row_shape.line.width = Pt(1)
        
        tf_obj = row_shape.text_frame
        tf_obj.word_wrap = True
        tf_obj.margin_left = Inches(0.2)
        tf_obj.margin_top = Inches(0.12)
        
        p_o = tf_obj.paragraphs[0]
        p_o.text = f"{i+1}.  {obj}"
        p_o.font.name = "Arial"
        p_o.font.size = Pt(15)
        p_o.font.color.rgb = COLOR_TEXT
        p_o.font.bold = True

    # ── Slide 6: System Architecture Flow ─────────────────────────────
    slide6 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide6)
    add_header(slide6, "Technical Design", "System Architecture Flow")
    
    # Description
    desc_box6 = slide6.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.6))
    tf_desc6 = desc_box6.text_frame
    tf_desc6.word_wrap = True
    p_d6 = tf_desc6.paragraphs[0]
    p_d6.text = "End-to-end data flow and architectural synchronization of components:"
    p_d6.font.name = "Arial"
    p_d6.font.size = Pt(16)
    p_d6.font.color.rgb = COLOR_MUTED
    
    # Flow Cards
    flow_steps = [
        {"num": "01", "label": "Chrome Extension", "desc": "Interactive DOM selector picker & AI validation"},
        {"num": "02", "label": "Scrapy / Playwright", "desc": "Routed extraction loops & anti-bot WAF bypass"},
        {"num": "03", "label": "FastAPI ETL Engine", "desc": "Deduplication, numeric parsing, & text engineering"},
        {"num": "04", "label": "ML Registry & APIs", "desc": "Auto-tuning models, scoring, & prediction API"}
    ]
    
    card_width = Inches(2.7)
    card_height = Inches(4.0)
    gap = Inches(0.3)
    arrow_width = Inches(0.3)
    
    for i, s in enumerate(flow_steps):
        # Position card
        left_card = Inches(0.8 + i * 3.0)
        top_card = Inches(2.4)
        
        # Draw step card
        card = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_card, top_card, card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD
        card.line.color.rgb = COLOR_BORDER
        card.line.width = Pt(1)
        
        tf_step = card.text_frame
        tf_step.word_wrap = True
        tf_step.margin_left = Inches(0.15)
        tf_step.margin_right = Inches(0.15)
        tf_step.margin_top = Inches(0.2)
        
        p_num = tf_step.paragraphs[0]
        p_num.text = s["num"]
        p_num.font.name = "Arial"
        p_num.font.size = Pt(20)
        p_num.font.bold = True
        p_num.font.color.rgb = COLOR_ACCENT
        
        p_lbl = tf_step.add_paragraph()
        p_lbl.text = "\n" + s["label"]
        p_lbl.font.name = "Arial"
        p_lbl.font.size = Pt(16)
        p_lbl.font.bold = True
        p_lbl.font.color.rgb = COLOR_TEXT
        
        p_desc = tf_step.add_paragraph()
        p_desc.text = "\n" + s["desc"]
        p_desc.font.name = "Arial"
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = COLOR_MUTED
        p_desc.alignment = PP_ALIGN.CENTER
        
        # Draw connecting arrow (if not last card)
        if i < 3:
            left_arrow = left_card + card_width + Inches(0.05)
            arrow = slide6.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left_arrow, top_card + Inches(1.8), Inches(0.2), Inches(0.4))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLOR_BORDER
            arrow.line.fill.background()

    # Save to disk
    output_filename = "presentation.pptx"
    prs.save(output_filename)
    print(f"Presentation successfully created: {os.path.abspath(output_filename)}")

if __name__ == "__main__":
    create_presentation()
